from pptx import Presentation
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from pptx.util import Inches, Pt

import os
import sys

prefix = sys.argv[1]
order = sys.argv[2]
dsize = sys.argv[3]
gpu1 = sys.argv[4]
gpu2 = sys.argv[5]
bw1 = sys.argv[6]
bw2 = sys.argv[7]
app = sys.argv[8]
path = sys.argv[9]

output_file = 'output.pptx'

# Check if the output file already exists
if os.path.exists(output_file):
    prs = Presentation(output_file)
else:
    # Create a new PowerPoint presentation
    prs = Presentation()

# Add a slide to the presentation
slide_layout = prs.slide_layouts[5]  # Layout for a title and content slide
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

# Add title to the slide
title_shape = slide.shapes.title
title_shape.text = prefix + "+" + app + "+" + order + "+" + dsize + "GB"

# Define the dimensions and positions of the PDFs on the slide
left = Inches(0.7)
top = Inches(2.5)
width = Inches(4)
height = Inches(2)
title_height = Inches(0.7)  # Height of the title text box

# pdf_files = ['cholesky_32_0.5_25.pdf', 'cholesky_32_0.5_25000000.pdf',
# 						 'cholesky_8_0.5_25.pdf', 'cholesky_8_0.5_25000000.pdf']

pdf_files = [f'{app}_{gpu1}_{dsize}_{bw1}.pdf', f'{app}_{gpu1}_{dsize}_{bw2}.pdf',
 						 f'{app}_{gpu2}_{dsize}_{bw1}.pdf', f'{app}_{gpu2}_{dsize}_{bw2}.pdf']

titles = [ f'{app} + {gpu1}GPUs + {dsize}GB + {bw1}GB/s', f'{app} + {gpu1}GPUs + {dsize}GB + {bw2}GB/s',
					 f'{app} + {gpu2}GPUs + {dsize}GB + {bw1}GB/s', f'{app} + {gpu2}GPUs + {dsize}GB + {bw2}GB/s' ]

# Insert images into the slide
for i, (pdf_file, title) in enumerate(zip(pdf_files, titles)):
    images = convert_from_path(path + "/" + pdf_file)
    image = images[0]  # Assuming you want the first page of each PDF
    image_stream = 'temp_image{}.png'.format(i)  # Temporary image file
    
    # Save image to a temporary file
    image.save(image_stream)
    
    # Insert image into the slide
    slide.shapes.add_picture(image_stream, left, top, width, height)

		# Insert each plot's text title into the slides
    title_shape = slide.shapes.add_textbox(left, top - title_height, width, title_height)
    title_text_frame = title_shape.text_frame
    p = title_text_frame.add_paragraph()
    p.text = title
    p.alignment = 1  # Center alignment

    print("text:", title)

    font = p.font
    font.size = Pt(13)  # Adjust the font size as desired
    
    if (i + 1) % 2 == 0:
        left = Inches(0.7)
        top = top + height + Inches(0.5)
    else:
        left = left + width + Inches(0.5)

# Save the PowerPoint presentation
prs.save('output.pptx')
